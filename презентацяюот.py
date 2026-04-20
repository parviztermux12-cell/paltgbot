import asyncio
import logging
import json
import os
import re
from datetime import datetime, date, timedelta
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, ChatMember, LabeledPrice
from telegram.ext import Application, CommandHandler, MessageHandler, CallbackQueryHandler, PreCheckoutQueryHandler, filters, ContextTypes
import aiohttp
from ddgs import DDGS
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image as PILImage
import zipfile

# ===== НАСТРОЙКИ =====
TELEGRAM_BOT_TOKEN = "8709660042:AAFqXswl8OwE1qewm4KMrQVgDMjV34iTro4"  # Замените на токен нового бота
PAYMENT_PROVIDER_TOKEN = ""  # Для Telegram Stars

CEREBRAS_API_KEYS = [
    "csk-fek5v5dn9cxj853hfk9cw3hvc24wwn3ddme63tmet8w96dmw",
    "csk-yh2rcf28e6tv9t9tfeynhd5xmfep8xcyc446h3tj3y5yc64j"
]

CEREBRAS_URL = "https://api.cerebras.ai/v1/chat/completions"
CEREBRAS_MODEL = "llama3.1-8b"

REQUIRED_CHANNEL_ID = -1003851572008
REQUIRED_CHANNEL_LINK = "https://t.me/izzzy_vpn"

# Лимиты
DAILY_PRESENTATION_LIMIT = 3  # для обычных пользователей
DAILY_PRESENTATION_LIMIT_VIP = 20  # для VIP
MAX_SLIDES_FREE = 12
MAX_SLIDES_VIP = 45

# VIP цены в Telegram Stars
VIP_PRICES = {"week": 25, "month": 50, "year": 250}

# Темы оформления
THEMES = {
    "business": {
        "name": "💼 Бизнес",
        "colors": ["#1a365d", "#2b6cb0", "#3182ce"],
        "bg_prompt": "professional business presentation background, blue tones, corporate style, clean design"
    },
    "creative": {
        "name": "🎨 Креатив",
        "colors": ["#702459", "#9b2c6d", "#d53f8c"],
        "bg_prompt": "creative colorful presentation background, artistic design, vibrant colors"
    },
    "tech": {
        "name": "🚀 Технологии",
        "colors": ["#1a202c", "#2d3748", "#4a5568"],
        "bg_prompt": "modern technology presentation background, digital, futuristic, dark theme with neon accents"
    },
    "nature": {
        "name": "🌿 Природа",
        "colors": ["#22543d", "#2f855a", "#48bb78"],
        "bg_prompt": "nature presentation background, green, eco-friendly, leaves, organic design"
    },
    "education": {
        "name": "📚 Образование",
        "colors": ["#2c5282", "#3182ce", "#63b3ed"],
        "bg_prompt": "educational presentation background, school theme, clean and professional"
    },
    "minimal": {
        "name": "⬜ Минимализм",
        "colors": ["#ffffff", "#e2e8f0", "#cbd5e0"],
        "bg_prompt": "minimalist white presentation background, clean, simple, elegant"
    }
}

logging.basicConfig(level=logging.INFO)

# ===== ФУНКЦИИ ДЛЯ РАБОТЫ С ДАННЫМИ =====
user_data_file = "preza_users.json"

def load_json(file):
    if os.path.exists(file):
        with open(file, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {}

def save_json(file, data):
    with open(file, 'w', encoding='utf-8') as f:
        json.dump(data, f, indent=2, ensure_ascii=False)

user_data = load_json(user_data_file)

def save_all():
    save_json(user_data_file, user_data)

def get_user(user_id):
    uid = str(user_id)
    if uid not in user_data:
        user_data[uid] = {
            "presentations_today": 0,
            "total_presentations": 0,
            "last_reset_date": str(date.today()),
            "vip_until": None,
            "referrer_id": None,
            "referrals": [],
            "bonus_presentations": 0
        }
    
    # Сброс ежедневного лимита
    if user_data[uid]["last_reset_date"] != str(date.today()):
        user_data[uid]["presentations_today"] = 0
        user_data[uid]["last_reset_date"] = str(date.today())
        save_all()
    
    # Проверка VIP
    if user_data[uid].get("vip_until"):
        vip_until = datetime.fromisoformat(user_data[uid]["vip_until"])
        if vip_until < datetime.now():
            user_data[uid]["vip_until"] = None
            save_all()
    
    return user_data[uid]

def get_daily_limit(user_id):
    user = get_user(user_id)
    if user.get("vip_until"):
        return DAILY_PRESENTATION_LIMIT_VIP + user.get("bonus_presentations", 0)
    return DAILY_PRESENTATION_LIMIT + user.get("bonus_presentations", 0)

def get_max_slides(user_id):
    user = get_user(user_id)
    return MAX_SLIDES_VIP if user.get("vip_until") else MAX_SLIDES_FREE

async def check_subscription(user_id, context):
    try:
        member = await context.bot.get_chat_member(REQUIRED_CHANNEL_ID, user_id)
        return member.status in [ChatMember.MEMBER, ChatMember.ADMINISTRATOR, ChatMember.OWNER]
    except:
        return False

# ===== ФУНКЦИИ ДЛЯ РАБОТЫ С AI =====
async def cerebras_request(messages, user_id=None, max_tokens=800, temperature=0.7):
    if user_id:
        user = get_user(user_id)
        start_index = user.get("current_key_index", 0)
    else:
        start_index = 0

    for attempt in range(len(CEREBRAS_API_KEYS)):
        key_index = (start_index + attempt) % len(CEREBRAS_API_KEYS)
        api_key = CEREBRAS_API_KEYS[key_index]
        try:
            async with aiohttp.ClientSession() as session:
                async with session.post(
                    CEREBRAS_URL,
                    headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                    json={"model": CEREBRAS_MODEL, "messages": messages, "max_tokens": max_tokens, "temperature": temperature},
                    timeout=aiohttp.ClientTimeout(total=120)
                ) as response:
                    if response.status == 200:
                        result = await response.json()
                        return result["choices"][0]["message"]["content"]
                    elif response.status == 429:
                        await asyncio.sleep(1)
                        continue
        except Exception as e:
            logging.error(f"Ошибка Cerebras: {e}")
            continue
    return None

async def generate_presentation_content(topic: str, slides_count: int, theme: str) -> dict:
    """Генерирует содержимое презентации"""
    theme_data = THEMES.get(theme, THEMES["minimal"])
    
    prompt = f"""Создай подробную презентацию на тему: "{topic}"
Количество слайдов: {slides_count}

Для каждого слайда укажи:
1. Заголовок
2. Содержание (3-5 предложений или пунктов)
3. Ключевые слова для поиска картинки (на английском, 3-5 слов)

Тема оформления: {theme_data['name']}

Формат ответа (строго JSON):
{{
    "title": "Название презентации",
    "slides": [
        {{
            "title": "Заголовок слайда",
            "content": "Текст слайда (3-5 предложений)",
            "image_keywords": "keywords for image search"
        }}
    ]
}}

Правила:
- Первый слайд - введение/титульный
- Последний слайд - заключение
- Содержание должно быть информативным и структурированным
- Используй профессиональный, деловой стиль"""

    messages = [
        {"role": "system", "content": "Ты профессиональный копирайтер и дизайнер презентаций. Отвечай только валидным JSON."},
        {"role": "user", "content": prompt}
    ]
    
    response = await cerebras_request(messages, max_tokens=4000, temperature=0.7)
    
    if response:
        try:
            # Очищаем ответ
            response = re.sub(r'```json\n?', '', response)
            response = re.sub(r'```\n?', '', response)
            return json.loads(response)
        except json.JSONDecodeError as e:
            logging.error(f"JSON ошибка: {e}")
            return None
    return None

async def search_image(keywords: str) -> str:
    """Ищет картинку по ключевым словам"""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.images(
                keywords,
                region="ru-ru",
                safesearch="moderate",
                max_results=1
            ))
            if results:
                return results[0].get("image")
    except Exception as e:
        logging.error(f"Ошибка поиска картинки: {e}")
    return None

async def download_image(url: str) -> bytes:
    """Скачивает изображение"""
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=aiohttp.ClientTimeout(total=15)) as response:
                if response.status == 200:
                    return await response.read()
    except Exception as e:
        logging.error(f"Ошибка скачивания: {e}")
    return None

# ===== ФУНКЦИИ ДЛЯ СОЗДАНИЯ PDF И PPTX =====
def create_pdf(presentation_data: dict, author: str, theme: str, images: list) -> BytesIO:
    """Создает PDF файл презентации"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    
    # Создаем стили
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        alignment=TA_CENTER,
        spaceAfter=30
    )
    
    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading2'],
        fontSize=18,
        alignment=TA_LEFT,
        spaceAfter=20,
        textColor='#1a365d'
    )
    
    content_style = ParagraphStyle(
        'Content',
        parent=styles['Normal'],
        fontSize=12,
        alignment=TA_LEFT,
        spaceAfter=10
    )
    
    story = []
    
    # Титульный слайд
    story.append(Paragraph(f"<b>{presentation_data['title']}</b>", title_style))
    story.append(Spacer(1, 50))
    story.append(Paragraph(f"Автор: {author}", content_style))
    story.append(Paragraph(f"Дата: {datetime.now().strftime('%d.%m.%Y')}", content_style))
    story.append(PageBreak())
    
    # Слайды
    for i, slide in enumerate(presentation_data['slides']):
        story.append(Paragraph(f"<b>{slide['title']}</b>", slide_title_style))
        story.append(Spacer(1, 10))
        
        # Добавляем картинку если есть
        if i < len(images) and images[i]:
            try:
                img_data = images[i]
                img = PILImage.open(BytesIO(img_data))
                img_path = BytesIO()
                img.save(img_path, format='PNG')
                img_path.seek(0)
                
                img_width = 150 * mm
                img_height = 100 * mm
                story.append(Image(img_path, width=img_width, height=img_height))
                story.append(Spacer(1, 10))
            except:
                pass
        
        # Добавляем текст
        content_text = slide['content'].replace('\n', '<br/>')
        story.append(Paragraph(content_text, content_style))
        story.append(PageBreak())
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def create_pptx(presentation_data: dict, author: str, images: list) -> BytesIO:
    """Создает PPTX файл презентации"""
    prs = PPTXPresentation()
    
    # Титульный слайд
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = presentation_data['title']
    subtitle.text = f"Автор: {author}\n{datetime.now().strftime('%d.%m.%Y')}"
    
    # Слайды с содержимым
    for i, slide_data in enumerate(presentation_data['slides']):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = slide_data['title']
        
        content = slide.placeholders[1]
        content.text = slide_data['content']
        
        # Добавляем картинку если есть
        if i < len(images) and images[i]:
            try:
                img_data = images[i]
                img_stream = BytesIO(img_data)
                slide.shapes.add_picture(img_stream, Inches(1), Inches(1), height=Inches(3))
            except:
                pass
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# ===== КОМАНДЫ И CALLBACK-И =====
async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик команды /start"""
    user = update.effective_user
    
    # Проверка подписки
    if not await check_subscription(user.id, context):
        keyboard = [[InlineKeyboardButton("📢 Подписаться на канал", url=REQUIRED_CHANNEL_LINK)]]
        await update.message.reply_text(
            f"✨ Добро пожаловать в Izzzy Preza ✨\n\n"
            f"Я бот для создания профессиональных презентаций!\n\n"
            f"❌ Для использования бота необходимо подписаться на наш канал.",
            reply_markup=InlineKeyboardMarkup(keyboard)
        )
        return
    
    # Приветствие
    await update.message.reply_text(
        f"✨ <b>Izzzy Preza</b> - твой персональный помощник для создания презентаций! ✨\n\n"
        f"📊 <b>Что я умею:</b>\n"
        f"• Создаю красивые презентации на любую тему\n"
        f"• Подбираю картинки для каждого слайда\n"
        f"• Экспортирую в PDF и PPTX форматы\n"
        f"• Поддерживаю до {MAX_SLIDES_FREE} слайдов (VIP: до {MAX_SLIDES_VIP})\n\n"
        f"🎯 <b>Как начать:</b>\n"
        f"Нажми на кнопку «➕ Создать презентацию» и следуй инструкциям!\n\n"
        f"⭐ Приобрети VIP для увеличения лимитов и приоритетной обработки!",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Создать презентацию", callback_data="create_presentation")],
            [InlineKeyboardButton("⭐ VIP статус", callback_data="vip_menu"), InlineKeyboardButton("👤 Мой профиль", callback_data="profile")],
            [InlineKeyboardButton("❓ Помощь", callback_data="help")]
        ])
    )

async def profile_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает профиль пользователя"""
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    user = get_user(user_id)
    
    daily_limit = get_daily_limit(user_id)
    remaining = daily_limit - user["presentations_today"]
    max_slides = get_max_slides(user_id)
    
    vip_status = "✅ Активен" if user.get("vip_until") else "❌ Неактивен"
    if user.get("vip_until"):
        vip_status += f"\n   до {datetime.fromisoformat(user['vip_until']).strftime('%d.%m.%Y')}"
    
    text = f"""👤 <b>Мой профиль</b>

📊 <b>Статистика:</b>
• Сегодня: {user['presentations_today']}/{daily_limit}
• Осталось: {remaining}
• Всего создано: {user['total_presentations']}

📄 <b>Лимиты:</b>
• Макс. слайдов: {max_slides}
• Форматы: PDF, PPTX

⭐ <b>VIP статус:</b> {vip_status}
👥 <b>Рефералов:</b> {len(user.get('referrals', []))}
🎁 <b>Бонусов:</b> {user.get('bonus_presentations', 0)}"""
    
    await query.edit_message_text(text, parse_mode="HTML", reply_markup=InlineKeyboardMarkup([
        [InlineKeyboardButton("‹ Назад", callback_data="back_to_menu")]
    ]))

async def help_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Помощь"""
    query = update.callback_query
    await query.answer()
    
    text = f"""❓ <b>Помощь и инструкция</b>

<b>Как создать презентацию:</b>
1. Нажми «➕ Создать презентацию»
2. Введи тему
3. Выбери количество слайдов
4. Выбери тему оформления
5. Укажи автора
6. Подтверди данные
7. Жди генерации (до 5 минут)

<b>Лимиты:</b>
• Обычный: {DAILY_PRESENTATION_LIMIT} презентаций/день, до {MAX_SLIDES_FREE} слайдов
• VIP: {DAILY_PRESENTATION_LIMIT_VIP} презентаций/день, до {MAX_SLIDES_VIP} слайдов

<b>Форматы вывода:</b>
• PDF - для печати и просмотра
• PPTX - для редактирования в PowerPoint

<b>Команды:</b>
/start - Главное меню
/help - Помощь

⭐ Приобрети VIP для увеличения лимитов!"""
    
    await query.edit_message_text(text, parse_mode="HTML", reply_markup=InlineKeyboardMarkup([
        [InlineKeyboardButton("‹ Назад", callback_data="back_to_menu")]
    ]))

async def back_to_menu(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Возврат в главное меню"""
    query = update.callback_query
    await query.answer()
    
    await query.edit_message_text(
        f"✨ <b>Izzzy Preza</b> - твой персональный помощник для создания презентаций! ✨\n\n"
        f"📊 <b>Что я умею:</b>\n"
        f"• Создаю красивые презентации на любую тему\n"
        f"• Подбираю картинки для каждого слайда\n"
        f"• Экспортирую в PDF и PPTX форматы\n\n"
        f"🎯 Нажми на кнопку «➕ Создать презентацию» чтобы начать!",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Создать презентацию", callback_data="create_presentation")],
            [InlineKeyboardButton("⭐ VIP статус", callback_data="vip_menu"), InlineKeyboardButton("👤 Мой профиль", callback_data="profile")],
            [InlineKeyboardButton("❓ Помощь", callback_data="help")]
        ])
    )

async def create_presentation_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Начало создания презентации"""
    query = update.callback_query
    await query.answer()
    user_id = query.from_user.id
    
    # Проверка лимита
    daily_limit = get_daily_limit(user_id)
    user = get_user(user_id)
    
    if user["presentations_today"] >= daily_limit:
        await query.edit_message_text(
            f"⚠️ <b>Лимит превышен!</b>\n\n"
            f"Сегодня вы использовали {user['presentations_today']}/{daily_limit} презентаций.\n\n"
            f"Купите VIP для увеличения лимита до {DAILY_PRESENTATION_LIMIT_VIP} презентаций в день!",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("⭐ Купить VIP", callback_data="vip_menu")],
                [InlineKeyboardButton("‹ Назад", callback_data="back_to_menu")]
            ])
        )
        return
    
    context.user_data["presentation_step"] = "topic"
    await query.edit_message_text(
        "📝 <b>Шаг 1/5: Введите тему презентации</b>\n\n"
        "Примеры:\n"
        "• Маркетинг в социальных сетях\n"
        "• Искусственный интеллект в бизнесе\n"
        "• Как открыть свой бизнес\n\n"
        "Отправьте тему одним сообщением:",
        parse_mode="HTML"
    )

async def handle_topic_input(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработка ввода темы"""
    if context.user_data.get("presentation_step") != "topic":
        return
    
    topic = update.message.text.strip()
    if len(topic) < 3:
        await update.message.reply_text("❌ Тема слишком короткая. Введите минимум 3 символа.")
        return
    
    context.user_data["presentation_topic"] = topic
    context.user_data["presentation_step"] = "slides_count"
    
    user_id = update.effective_user.id
    max_slides = get_max_slides(user_id)
    
    # Создаем кнопки для выбора количества слайдов
    buttons = []
    for i in [5, 8, 10, 12]:
        if i <= max_slides:
            buttons.append([InlineKeyboardButton(f"{i} слайдов", callback_data=f"slides_{i}")])
    
    if max_slides > 12:
        buttons.append([InlineKeyboardButton(f"15 слайдов", callback_data="slides_15")])
        buttons.append([InlineKeyboardButton(f"20 слайдов", callback_data="slides_20")])
        buttons.append([InlineKeyboardButton(f"30 слайдов", callback_data="slides_30")])
        buttons.append([InlineKeyboardButton(f"45 слайдов", callback_data="slides_45")])
    
    buttons.append([InlineKeyboardButton("‹ Назад", callback_data="back_to_slides")])
    
    await update.message.reply_text(
        f"📊 <b>Шаг 2/5: Выберите количество слайдов</b>\n\n"
        f"Тема: {topic}\n"
        f"Максимум: {max_slides} слайдов\n\n"
        f"Выберите из предложенных вариантов:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(buttons)
    )

async def slides_count_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Выбор количества слайдов"""
    query = update.callback_query
    await query.answer()
    
    if query.data == "back_to_slides":
        context.user_data["presentation_step"] = "topic"
        await query.edit_message_text(
            "📝 <b>Шаг 1/5: Введите тему презентации</b>\n\n"
            "Отправьте тему одним сообщением:",
            parse_mode="HTML"
        )
        return
    
    slides_count = int(query.data.split("_")[1])
    context.user_data["presentation_slides"] = slides_count
    context.user_data["presentation_step"] = "theme"
    
    # Показываем выбор темы оформления
    keyboard = []
    for theme_id, theme_data in THEMES.items():
        keyboard.append([InlineKeyboardButton(theme_data["name"], callback_data=f"theme_{theme_id}")])
    keyboard.append([InlineKeyboardButton("‹ Назад", callback_data="back_to_slides_count")])
    
    await query.edit_message_text(
        f"🎨 <b>Шаг 3/5: Выберите тему оформления</b>\n\n"
        f"Тема: {context.user_data['presentation_topic']}\n"
        f"Слайдов: {slides_count}\n\n"
        f"Как должен выглядеть фон презентации?",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup(keyboard)
    )

async def theme_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Выбор темы оформления"""
    query = update.callback_query
    await query.answer()
    
    if query.data == "back_to_slides_count":
        user_id = query.from_user.id
        max_slides = get_max_slides(user_id)
        
        buttons = []
        for i in [5, 8, 10, 12]:
            if i <= max_slides:
                buttons.append([InlineKeyboardButton(f"{i} слайдов", callback_data=f"slides_{i}")])
        
        if max_slides > 12:
            buttons.append([InlineKeyboardButton(f"15 слайдов", callback_data="slides_15")])
            buttons.append([InlineKeyboardButton(f"20 слайдов", callback_data="slides_20")])
            buttons.append([InlineKeyboardButton(f"30 слайдов", callback_data="slides_30")])
            buttons.append([InlineKeyboardButton(f"45 слайдов", callback_data="slides_45")])
        
        buttons.append([InlineKeyboardButton("‹ Назад", callback_data="back_to_slides")])
        
        await query.edit_message_text(
            f"📊 <b>Шаг 2/5: Выберите количество слайдов</b>\n\n"
            f"Тема: {context.user_data['presentation_topic']}\n"
            f"Максимум: {max_slides} слайдов",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup(buttons)
        )
        return
    
    theme_id = query.data.split("_")[1]
    context.user_data["presentation_theme"] = theme_id
    context.user_data["presentation_step"] = "author"
    
    # Показываем пример оформления
    theme_data = THEMES[theme_id]
    
    await query.edit_message_text(
        f"✍️ <b>Шаг 4/5: Укажите автора презентации</b>\n\n"
        f"📌 Тема: {context.user_data['presentation_topic']}\n"
        f"📄 Слайдов: {context.user_data['presentation_slides']}\n"
        f"🎨 Оформление: {theme_data['name']}\n\n"
        f"<b>Пример оформления:</b>\n"
        f"Фон будет в стиле: {theme_data['bg_prompt'][:100]}...\n\n"
        f"Введите ФИО автора:",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("‹ Назад к темам", callback_data="back_to_themes")]
        ])
    )

async def author_input(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Ввод автора"""
    if context.user_data.get("presentation_step") != "author":
        return
    
    author = update.message.text.strip()
    if len(author) < 2:
        await update.message.reply_text("❌ Введите корректное ФИО (минимум 2 символа)")
        return
    
    context.user_data["presentation_author"] = author
    context.user_data["presentation_step"] = "confirm_author"
    
    await update.message.reply_text(
        f"✅ <b>Подтвердите данные автора</b>\n\n"
        f"ФИО: <b>{author}</b>\n\n"
        f"Всё верно?",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("✅ Да, верно", callback_data="confirm_author_yes"),
             InlineKeyboardButton("❌ Нет, исправить", callback_data="confirm_author_no")]
        ])
    )

async def confirm_author_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Подтверждение автора"""
    query = update.callback_query
    await query.answer()
    
    if query.data == "confirm_author_no":
        context.user_data["presentation_step"] = "author"
        await query.edit_message_text(
            "✍️ <b>Введите ФИО автора заново:</b>",
            parse_mode="HTML"
        )
        return
    
    # Подтверждено, переходим к подтверждению темы
    context.user_data["presentation_step"] = "confirm_topic"
    
    await query.edit_message_text(
        f"📝 <b>Подтвердите тему презентации</b>\n\n"
        f"Тема: <b>{context.user_data['presentation_topic']}</b>\n\n"
        f"Всё верно?",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("✅ Да, всё верно", callback_data="confirm_topic_yes"),
             InlineKeyboardButton("❌ Нет, исправить", callback_data="confirm_topic_no")]
        ])
    )

async def confirm_topic_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Подтверждение темы"""
    query = update.callback_query
    await query.answer()
    
    if query.data == "confirm_topic_no":
        context.user_data["presentation_step"] = "topic"
        await query.edit_message_text(
            "📝 <b>Введите тему презентации заново:</b>",
            parse_mode="HTML"
        )
        return
    
    # Всё подтверждено, начинаем генерацию
    await query.edit_message_text(
        "🍫 <b>Начинаю создавать презентацию...</b>\n\n"
        "📝 Пишу текст и подбираю фотографии под тематику.\n"
        "⏳ Генерация может идти до 5 минут...\n\n"
        "Пожалуйста, подождите, я стараюсь сделать всё красиво! 🎨",
        parse_mode="HTML"
    )
    
    # Запускаем генерацию
    user_id = query.from_user.id
    user = get_user(user_id)
    
    # Генерируем содержимое
    presentation = await generate_presentation_content(
        context.user_data["presentation_topic"],
        context.user_data["presentation_slides"],
        context.user_data["presentation_theme"]
    )
    
    if not presentation:
        await query.edit_message_text(
            "❌ <b>Ошибка генерации!</b>\n\n"
            "Не удалось создать презентацию. Попробуйте позже или измените тему.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton("‹ Назад в меню", callback_data="back_to_menu")]
            ])
        )
        return
    
    # Подбираем картинки
    status_msg = await query.message.edit_text(
        "🎨 <b>Создаю слайды и подбираю иллюстрации...</b>\n\n"
        "Это может занять ещё 1-2 минуты...",
        parse_mode="HTML"
    )
    
    images = []
    for i, slide in enumerate(presentation['slides']):
        await status_msg.edit_text(
            f"🎨 <b>Создаю презентацию...</b>\n\n"
            f"Слайд {i+1}/{len(presentation['slides'])}: {slide['title'][:30]}...\n\n"
            f"⏳ Ищу подходящую картинку...",
            parse_mode="HTML"
        )
        
        img_url = await search_image(slide.get('image_keywords', slide['title']))
        if img_url:
            img_data = await download_image(img_url)
            if img_data:
                images.append(img_data)
            else:
                images.append(None)
        else:
            images.append(None)
    
    # Создаем файлы
    await status_msg.edit_text("📄 <b>Создаю PDF и PPTX файлы...</b>", parse_mode="HTML")
    
    pdf_buffer = create_pdf(presentation, context.user_data["presentation_author"], context.user_data["presentation_theme"], images)
    pptx_buffer = create_pptx(presentation, context.user_data["presentation_author"], images)
    
    # Обновляем статистику
    user["presentations_today"] += 1
    user["total_presentations"] += 1
    save_all()
    
    # Отправляем файлы
    await status_msg.delete()
    
    await update.callback_query.message.reply_document(
        document=pdf_buffer,
        filename=f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
        caption=f"✅ <b>Презентация готова!</b>\n\n"
                f"📌 <b>Название:</b> {presentation['title']}\n"
                f"📄 <b>Слайдов:</b> {len(presentation['slides'])}\n"
                f"👤 <b>Автор:</b> {context.user_data['presentation_author']}\n\n"
                f"Файлы готовы к скачиванию! 📥",
        parse_mode="HTML"
    )
    
    await update.callback_query.message.reply_document(
        document=pptx_buffer,
        filename=f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
        caption="📊 <b>Презентация в формате PPTX</b>\n\nМожно редактировать в PowerPoint! 🎨",
        parse_mode="HTML"
    )
    
    # Очищаем данные
    context.user_data.clear()

# ===== VIP ФУНКЦИИ =====
async def vip_menu_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Меню VIP"""
    query = update.callback_query
    await query.answer()
    user = get_user(query.from_user.id)
    
    status = f"✅ Активен до {datetime.fromisoformat(user['vip_until']).strftime('%d.%m.%Y')}" if user.get("vip_until") else "❌ Неактивен"
    
    text = f"""⭐ <b>VIP статус Izzzy Preza</b>

┌ 📊 Текущий статус: {status}
├ 🎁 Бонусы VIP:
├   • {DAILY_PRESENTATION_LIMIT_VIP} презентаций/день
├   • До {MAX_SLIDES_VIP} слайдов
├   • Приоритетная обработка
├   • Эксклюзивные темы оформления
└ 🚀 Без рекламы и ограничений

<b>💎 Цены (Telegram Stars):</b>
• 7 дней — {VIP_PRICES['week']} ⭐
• 30 дней — {VIP_PRICES['month']} ⭐
• 365 дней — {VIP_PRICES['year']} ⭐

Нажми на кнопку ниже для покупки:"""
    
    await query.edit_message_text(
        text,
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton(f"7 дней — {VIP_PRICES['week']} ⭐", callback_data="vip_week")],
            [InlineKeyboardButton(f"30 дней — {VIP_PRICES['month']} ⭐", callback_data="vip_month")],
            [InlineKeyboardButton(f"365 дней — {VIP_PRICES['year']} ⭐", callback_data="vip_year")],
            [InlineKeyboardButton("‹ Назад", callback_data="back_to_menu")]
        ])
    )

async def vip_purchase_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Покупка VIP"""
    query = update.callback_query
    await query.answer()
    duration = query.data.split("_")[1]
    price = VIP_PRICES[duration]
    
    await context.bot.send_invoice(
        chat_id=query.message.chat_id,
        title=f"VIP {duration} Izzzy Preza",
        description=f"До {MAX_SLIDES_VIP} слайдов, {DAILY_PRESENTATION_LIMIT_VIP} презентаций/день",
        payload=f"vip_{duration}_{query.from_user.id}",
        provider_token=PAYMENT_PROVIDER_TOKEN,
        currency="XTR",
        prices=[LabeledPrice("VIP доступ", price)]
    )

async def pre_checkout_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Pre-checkout"""
    await update.pre_checkout_query.answer(ok=True)

async def successful_payment_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Успешная оплата"""
    msg = update.message
    user_id = msg.from_user.id
    duration = msg.successful_payment.invoice_payload.split("_")[1]
    delta = {"week": timedelta(days=7), "month": timedelta(days=30), "year": timedelta(days=365)}[duration]
    
    user = get_user(user_id)
    user["vip_until"] = (datetime.now() + delta).isoformat()
    save_all()
    
    await msg.reply_text(
        f"✅ <b>VIP активирован на {duration}!</b>\n\n"
        f"Теперь вам доступно:\n"
        f"• {DAILY_PRESENTATION_LIMIT_VIP} презентаций в день\n"
        f"• До {MAX_SLIDES_VIP} слайдов\n"
        f"• Приоритетная обработка\n\n"
        f"Спасибо за поддержку! 🎉",
        parse_mode="HTML",
        reply_markup=InlineKeyboardMarkup([
            [InlineKeyboardButton("➕ Создать презентацию", callback_data="create_presentation")]
        ])
    )

# ===== ОСНОВНОЙ ЗАПУСК =====
def main():
    app = Application.builder().token(TELEGRAM_BOT_TOKEN).build()
    
    # Command handlers
    app.add_handler(CommandHandler("start", start_command))
    app.add_handler(CommandHandler("help", help_callback))
    
    # Callback handlers
    app.add_handler(CallbackQueryHandler(back_to_menu, pattern="back_to_menu"))
    app.add_handler(CallbackQueryHandler(profile_callback, pattern="profile"))
    app.add_handler(CallbackQueryHandler(help_callback, pattern="help"))
    app.add_handler(CallbackQueryHandler(create_presentation_start, pattern="create_presentation"))
    app.add_handler(CallbackQueryHandler(slides_count_callback, pattern="^slides_\\d+$"))
    app.add_handler(CallbackQueryHandler(slides_count_callback, pattern="back_to_slides"))
    app.add_handler(CallbackQueryHandler(theme_callback, pattern="^theme_"))
    app.add_handler(CallbackQueryHandler(theme_callback, pattern="back_to_themes"))
    app.add_handler(CallbackQueryHandler(confirm_author_callback, pattern="^confirm_author_"))
    app.add_handler(CallbackQueryHandler(confirm_topic_callback, pattern="^confirm_topic_"))
    app.add_handler(CallbackQueryHandler(vip_menu_callback, pattern="vip_menu"))
    app.add_handler(CallbackQueryHandler(vip_purchase_callback, pattern="vip_week|vip_month|vip_year"))
    
    # Payment handlers
    app.add_handler(PreCheckoutQueryHandler(pre_checkout_callback))
    app.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, successful_payment_callback))
    
    # Message handlers
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_topic_input))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, author_input))
    
    print("✅ Бот Izzzy Preza запущен!")
    app.run_polling()

if __name__ == "__main__":
    main()